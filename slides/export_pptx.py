#!/usr/bin/env python3
"""Export WP4.4 slides to PPTX for Google Slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0F, 0x1B, 0x2D)
SLATE = RGBColor(0x33, 0x41, 0x55)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)

DIR = os.path.dirname(os.path.abspath(__file__))


def img(name):
    return os.path.join(DIR, name)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                bold=False, color=SLATE, alignment=PP_ALIGN.LEFT,
                font_name="Inter", italic=False, spacing_after=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = spacing_after
    return tf


def add_para(tf, text, font_size=14, bold=False, color=SLATE,
             alignment=PP_ALIGN.LEFT, font_name="Inter", italic=False,
             spacing_after=Pt(4), spacing_before=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = spacing_after
    p.space_before = spacing_before
    return p


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            return slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(path, left, top, height=height)
        else:
            return slide.shapes.add_picture(path, left, top)
    return None


def add_bg_image(slide, path, opacity=1.0):
    """Add a full-slide background image with optional opacity (0.0-1.0)."""
    pic = add_image_safe(slide, path, Inches(0), Inches(0),
                         width=SLIDE_W, height=SLIDE_H)
    if pic and opacity < 1.0:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        blip = pic._element.findall('.//a:blip', nsmap)[0]
        alpha_val = int(opacity * 100000)  # in 1000ths of a percent
        alpha_el = etree.SubElement(blip, '{http://schemas.openxmlformats.org/drawingml/2006/main}alphaModFix')
        alpha_el.set('amt', str(alpha_val))
    return pic


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ── 1. TITLE ──────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, NAVY)
add_bg_image(slide, img("title-bg.png"), opacity=0.6)

tf = add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.5),
            "Designing Quantum Interfaces", font_size=36, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER, font_name="Libre Baskerville")
add_para(tf, "for Society", font_size=36, bold=True, color=WHITE,
         alignment=PP_ALIGN.CENTER, font_name="Libre Baskerville", spacing_before=Pt(0))

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.5),
            "WP4.4 — Mental Models for Quantum Computing", font_size=18,
            color=RGBColor(0xA0, 0xB0, 0xCC), alignment=PP_ALIGN.CENTER)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(5.9), Inches(3.5), Inches(1.5), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

tf = add_textbox(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(1.5),
                 "", font_size=16, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
tf.paragraphs[0].text = "J. Derek Lomas — Assistant Professor of Positive AI"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.name = "Inter"
add_para(tf, "Caiseal Beardow — PhD Candidate in Human-Centered Design",
         font_size=16, color=WHITE,
         alignment=PP_ALIGN.CENTER)
add_para(tf, "Industrial Design Engineering & QuTech, TU Delft",
         font_size=13, color=RGBColor(0xBB, 0xCC, 0xDD),
         alignment=PP_ALIGN.CENTER, spacing_before=Pt(8))
add_para(tf, "Quantum Inspire Program Review — February 2026",
         font_size=11, color=RGBColor(0x88, 0x99, 0xAA),
         alignment=PP_ALIGN.CENTER, spacing_before=Pt(20))


# ── 2. THE PROGRAM ───────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "The Program", font_size=28, bold=True, color=NAVY, font_name="Libre Baskerville")

add_image_safe(slide, img("wp-diagram.png"),
               Inches(2.5), Inches(1.2), width=Inches(8))

tf = add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.5),
                 "", font_size=14, color=MUTED, italic=True)
tf.paragraphs[0].text = (
    '"The design of the Quantum Inspire interface is repeatedly adapted to '
    'incorporate the findings from stakeholders (WP5-7) and industrial design '
    'researchers, who investigate intuitive mental models of quantum bits and computers."'
)
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = MUTED


# ── 3. MANDATE ────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "WP4.4 Mandate", font_size=28, bold=True, color=NAVY, font_name="Libre Baskerville")

tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1),
                 "", font_size=14, italic=True, color=MUTED)
tf.paragraphs[0].text = (
    '"WP4.4 requires fundamental research on mental models for quantum '
    'computing concepts which will guide the UX/UI design."')
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = MUTED

tf2 = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2),
                  "", font_size=14, italic=True, color=MUTED)
tf2.paragraphs[0].text = (
    '"Deploy qualitative methods to create a common language (and shared '
    'mental models) acceptable to both experts and non-experts... technology '
    'probes, research-through-design, experience prototyping, storytelling and co-creation."')
tf2.paragraphs[0].font.size = Pt(14)
tf2.paragraphs[0].font.italic = True
tf2.paragraphs[0].font.color.rgb = MUTED

add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.4),
            "PROMISED DELIVERABLES", font_size=12, bold=True, color=BLUE)

bullets = [
    '"PhD thesis (month 48)" + "3 scientific publications (month 18, 30, 42)"',
    '"Updated user experience and interface of quantum-inspire.com (month 6, month 36)"',
    '"Data... available through public repositories such as 4TU or Zenodo"',
]
tf3 = add_textbox(slide, Inches(1), Inches(4.1), Inches(11), Inches(2),
                  "", font_size=14, color=SLATE)
for i, b in enumerate(bullets):
    if i == 0:
        tf3.paragraphs[0].text = f"•  {b}"
        tf3.paragraphs[0].font.size = Pt(14)
        tf3.paragraphs[0].font.color.rgb = SLATE
    else:
        add_para(tf3, f"•  {b}", font_size=14, color=SLATE, spacing_before=Pt(4))


# ── 4. THREE LINES OF WORK ───────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "Three Lines of Work", font_size=28, bold=True, color=NAVY, font_name="Libre Baskerville")

# Three boxes
box_data = [
    ("1. PhD Thesis", "Quantum Computational Thinking — framework, metaphors, and prototype for computing-focused QC education", GREEN),
    ("2. MSc Graduation", 'D\'Arcangelis (2024) — "How Might ChatGPT Improve the Accessibility of Quantum Computing?"', GREEN),
    ("3. Haiqu: AI Agent", "AI-mediated quantum computing — replication study + benchmark + interactive platform", BLUE),
]

for i, (title, desc, col) in enumerate(box_data):
    left = Inches(0.8 + i * 4.0)
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(1.3), Inches(3.6), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4) if col == GREEN else SURFACE
    box.line.color.rgb = col
    box.line.width = Pt(1.5)

    add_textbox(slide, left + Inches(0.2), Inches(1.5), Inches(3.2), Inches(0.4),
                title, font_size=14, bold=True, color=col)
    add_textbox(slide, left + Inches(0.2), Inches(2.0), Inches(3.2), Inches(1.2),
                desc, font_size=12, color=SLATE)

tf = add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1),
                 "", font_size=14, italic=True, color=MUTED)
tf.paragraphs[0].text = (
    '"Translate fundamental research in quantum physics and design to '
    'concrete outcomes for society via an accessible interface for quantum computers."')
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = MUTED


# ── 5. TRANSITION: PhD ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, NAVY)
add_bg_image(slide, img("grid-illustration.png"), opacity=0.25)

tf = add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
            "What is quantum computational",
            font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
            font_name="Libre Baskerville")
add_para(tf, "thinking and how do we teach it?",
         font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
         font_name="Libre Baskerville", spacing_before=Pt(0))

add_textbox(slide, Inches(1.5), Inches(4.0), Inches(8), Inches(0.5),
            "PhD RESEARCH — CAISEAL BEARDOW", font_size=12,
            color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.LEFT)


# ── 6. PhD (1/5): QCT ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "PhD (1/5): Quantum Computational Thinking", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

# Left: concepts table
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(3), Inches(0.3),
            "7 KEY QC CONCEPTS", font_size=11, bold=True, color=BLUE)

concepts = [
    ("Superposition", "Low"), ("Entanglement", "Low"), ("Probability", "Low"),
    ("Measurement", "Medium"), ("Qubits", "Medium"),
    ("Gate operations", "High"), ("Algorithms", "High"),
]

from pptx.util import Inches as In
tbl = slide.shapes.add_table(len(concepts)+1, 2, Inches(0.8), Inches(1.4), Inches(4), Inches(2.8))
table = tbl.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(1.5)

# Header
for ci, txt in enumerate(["Concept", "Complexity"]):
    cell = table.cell(0, ci)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY

for ri, (c, cx) in enumerate(concepts):
    for ci, val in enumerate([c, cx]):
        cell = table.cell(ri+1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = SLATE
        if ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE

# Right: CT dimensions
add_textbox(slide, Inches(5.5), Inches(1.0), Inches(5), Inches(0.3),
            "5 COMPUTATIONAL THINKING DIMENSIONS", font_size=11, bold=True, color=BLUE)

dims = [
    "Abstraction — categorising, generalising, relating",
    "Algorithmic thinking — decomposing, iterating, planning",
    "Data — types, values, transforming",
    "Logic — deducing, extrapolating, predicting",
    "Problem-solving — adjusting, judging",
]
tf = add_textbox(slide, Inches(5.5), Inches(1.4), Inches(7), Inches(2.5),
                 "", font_size=13, color=SLATE)
for i, d in enumerate(dims):
    if i == 0:
        tf.paragraphs[0].text = d
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = SLATE
    else:
        add_para(tf, d, font_size=13, color=SLATE, spacing_before=Pt(6))

# Key insight box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.5), Inches(4.2), Inches(7), Inches(1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
box.line.color.rgb = AMBER
box.line.width = Pt(1)

add_textbox(slide, Inches(5.7), Inches(4.35), Inches(6.5), Inches(0.8),
            "Key insight: Educators frame QC bottom-up from physics. Computing learners need the inverse — top-down from algorithms.",
            font_size=12, color=SLATE)


# ── 6. PhD (2/5): METAPHORS ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "PhD (2/5): Metaphor Research", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

# Left: existing metaphors
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.3),
            "WHAT EDUCATORS USE TODAY", font_size=11, bold=True, color=BLUE)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
            "39 metaphors collected via interviews with QC educators:",
            font_size=12, color=SLATE)

met_existing = [
    ("Algorithms", "Incremental construction"),
    ("Entanglement", "Paired objects"),
    ("Qubits", "A polar object"),
    ("Superposition", "Expansion of resources"),
]
tbl = slide.shapes.add_table(len(met_existing)+1, 2, Inches(0.8), Inches(1.9), Inches(5), Inches(1.6))
table = tbl.table
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(3)
for ci, txt in enumerate(["Concept", "Prevalent model"]):
    cell = table.cell(0, ci)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
for ri, (c, m) in enumerate(met_existing):
    for ci, val in enumerate([c, m]):
        cell = table.cell(ri+1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = SLATE

# Warn box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.7), Inches(5), Inches(0.9))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
box.line.color.rgb = AMBER
box.line.width = Pt(1)
add_textbox(slide, Inches(1), Inches(3.8), Inches(4.6), Inches(0.7),
            "Scored well on explainable (4.0) and relatable (3.8) but poorly on actionable (2.6/5). They describe properties, not computations.",
            font_size=11, color=SLATE)

# Right: new metaphors
add_textbox(slide, Inches(6.5), Inches(1.0), Inches(6), Inches(0.3),
            "WHAT WE DESIGNED", font_size=11, bold=True, color=BLUE)

met_new = [
    ("Qubits", "Coloured cells in a grid"),
    ("Gates", "Formatting cells"),
    ("Superposition", "Colour gradient on a cell"),
    ("Entanglement", "Opacity of adjacent cells"),
    ("Measurement", "Setting colour to a point"),
    ("Probability", "Manipulating colour in a gradient"),
    ("Algorithms", "Refining a route"),
]
tbl2 = slide.shapes.add_table(len(met_new)+1, 2, Inches(6.5), Inches(1.4), Inches(6), Inches(2.8))
table2 = tbl2.table
table2.columns[0].width = Inches(2.2)
table2.columns[1].width = Inches(3.8)
for ci, txt in enumerate(["Concept", "New metaphor"]):
    cell = table2.cell(0, ci)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
for ri, (c, m) in enumerate(met_new):
    for ci, val in enumerate([c, m]):
        cell = table2.cell(ri+1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = SLATE


# ── 7. PhD (3/5): DESIGN ITERATIONS ──────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "PhD (3/5): Design Iterations", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
            "From concept to functional prototype through iterative research-through-design:",
            font_size=14, color=SLATE)

add_image_safe(slide, img("phd-design-iterations.png"),
               Inches(1.5), Inches(1.6), height=Inches(4.4))

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
            "Early concept sketches, physical prototyping, and digital interface exploration",
            font_size=11, italic=True, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── 8. PhD (4/5): STORYBOARD & PROTOTYPE ─────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            'PhD (4/5): Prototype — "Guess a Cell"', font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

# Storyboards side by side
add_image_safe(slide, img("phd-storyboard-1.png"),
               Inches(0.5), Inches(1.0), width=Inches(4.5))
add_image_safe(slide, img("phd-storyboard-2.png"),
               Inches(5.2), Inches(1.0), width=Inches(4.5))

add_textbox(slide, Inches(0.5), Inches(4.6), Inches(9.2), Inches(0.4),
            'Storyboard: learner navigates rounds of "find the most red cell", scores improve, landscape reveals',
            font_size=11, italic=True, color=MUTED, alignment=PP_ALIGN.CENTER)

# Annotated prototype
add_image_safe(slide, img("phd-prototype-annotated.png"),
               Inches(9.8), Inches(1.0), height=Inches(3.8))

add_textbox(slide, Inches(9.8), Inches(5.0), Inches(3.2), Inches(0.5),
            "Annotated interface:\ncontrols, grid, mask layer",
            font_size=10, italic=True, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── 9. PhD (5/5): RESULTS ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "PhD (5/5): Evaluation Results", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

# Left: study info
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(3), Inches(0.3),
            "STUDY DESIGN", font_size=11, bold=True, color=BLUE)

study_points = [
    "•  10 participants with computing background, no QC knowledge",
    "•  Think-aloud protocol",
    "•  Deductive coding against QCT framework",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(1.5),
                 "", font_size=12, color=SLATE)
for i, pt in enumerate(study_points):
    if i == 0:
        tf.paragraphs[0].text = pt
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = SLATE
    else:
        add_para(tf, pt, font_size=12, color=SLATE)

# Result boxes
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.0), Inches(5.5), Inches(0.6))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
box.line.color.rgb = GREEN
box.line.width = Pt(1)
add_textbox(slide, Inches(1), Inches(3.05), Inches(5), Inches(0.5),
            "All 7 QC concepts and all 5 CT dimensions identified in every transcript.",
            font_size=12, bold=True, color=SLATE)

add_textbox(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.3),
            "CT DIMENSION BREAKDOWN", font_size=11, bold=True, color=BLUE)

add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
            "Abstraction 28%  ·  Data 21%  ·  Algorithmic 19%  ·  Problem-solving 18%  ·  Logic 14%",
            font_size=12, color=SLATE)

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(4.9), Inches(5.5), Inches(0.6))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
box2.line.color.rgb = GREEN
box2.line.width = Pt(1)
add_textbox(slide, Inches(1), Inches(4.95), Inches(5), Inches(0.5),
            "Participants connected multiple QC concepts through the shared metaphor system without explicit instruction.",
            font_size=11, color=SLATE)

# Right: results table image
add_image_safe(slide, img("phd-ct-results-table.png"),
               Inches(6.8), Inches(1.0), height=Inches(4.8))

add_textbox(slide, Inches(6.8), Inches(5.9), Inches(6), Inches(0.4),
            "CT dimension percentages per participant (Table 7.7)",
            font_size=10, italic=True, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── 11. TRANSITION: AI ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, NAVY)

tf = add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
            "What if AI could be",
            font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
            font_name="Libre Baskerville")
add_para(tf, "the quantum interface?",
         font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
         font_name="Libre Baskerville", spacing_before=Pt(0))

add_textbox(slide, Inches(1.5), Inches(4.0), Inches(8), Inches(0.5),
            "FROM METAPHORS TO AGENTS", font_size=12,
            color=RGBColor(0x66, 0x77, 0x88), alignment=PP_ALIGN.LEFT)


# ── 12. HAIQU.ORG ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
            "haiqu.org — AI Meets Quantum Hardware", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
            "haiqu.org is the interface between people and quantum computers — ask a question in plain language, get results from real hardware.",
            font_size=14, color=SLATE)

imgs_web = ["website-home-1200.png", "experiments.png", "explore.png"]
for i, im in enumerate(imgs_web):
    add_image_safe(slide, img(im),
                   Inches(0.5 + i * 4.2), Inches(1.6), width=Inches(3.9))

tf = add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
            "Interactive dashboard — experiments, results, visualisations.",
            font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)
add_para(tf, "All built by AI agents through natural language.",
         font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER, spacing_before=Pt(0))


# ── 13. MSc ──────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "MSc: Can ChatGPT Make QC Accessible?", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(7), Inches(0.4),
            "C. D'Arcangelis (2024, supervised with Deborah Nas)", font_size=14, bold=True, color=SLATE)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.3),
            "RESEARCH QUESTION", font_size=11, bold=True, color=BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(0.4),
            "How might ChatGPT improve the accessibility of quantum computing?",
            font_size=14, color=SLATE)

add_textbox(slide, Inches(0.8), Inches(2.4), Inches(3), Inches(0.3),
            "FINDINGS", font_size=11, bold=True, color=BLUE)

findings = [
    "•  ChatGPT as effective intermediary for democratising QC access",
    "•  Benefits at all skill levels: beginners do creative coding, programmers gain QM insight",
    '•  Created "Quantum Buddy 2.0" — custom GPT refined for non-experts',
    "•  Non-experts solved encryption problems with no QM background",
]
tf = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(7), Inches(2),
                 "", font_size=13, color=SLATE)
for i, f in enumerate(findings):
    if i == 0:
        tf.paragraphs[0].text = f
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = SLATE
    else:
        add_para(tf, f, font_size=13, color=SLATE, spacing_before=Pt(4))

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(4.8), Inches(7), Inches(0.7))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
box.line.color.rgb = GREEN
box.line.width = Pt(1)
add_textbox(slide, Inches(1), Inches(4.85), Inches(6.5), Inches(0.6),
            "Early evidence that LLMs can function as an accessibility layer for quantum computing.",
            font_size=13, bold=True, color=SLATE)

# Right side: AI bridge image
add_image_safe(slide, img("ai-bridge.png"),
               Inches(8), Inches(1.5), width=Inches(5))

add_textbox(slide, Inches(8), Inches(5.6), Inches(5), Inches(0.5),
            "AI as the interface between humans and quantum hardware",
            font_size=11, italic=True, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── 11. HAIQU: QUESTION ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "Haiqu: How Good Is AI at Quantum?", font_size=26, bold=True, color=NAVY, font_name="Libre Baskerville")

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.3),
            "QISKIT HUMANEVAL (IBM, 2024)", font_size=11, bold=True, color=BLUE)
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.6),
            "151 quantum programming tasks. LLMs achieve 62–71% zero-shot, 79.5% with retrieval augmentation.",
            font_size=13, color=SLATE)

add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5), Inches(0.3),
            "OUR REPLICATION STUDY", font_size=11, bold=True, color=BLUE)
add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.6),
            "Go beyond code generation: can an AI agent replicate published experiments on real hardware, end-to-end?",
            font_size=13, color=SLATE, italic=True)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.2))
box.fill.solid()
box.fill.fore_color.rgb = SURFACE
box.line.color.rgb = BLUE
box.line.width = Pt(1)
add_textbox(slide, Inches(1), Inches(3.4), Inches(5), Inches(1),
            "Replication is some of the easiest science to do — the protocol exists, the claims are defined. If AI can do this, it validates the approach. In the hands of experts, the same tools could accelerate greenfield science.",
            font_size=12, color=SLATE)

add_image_safe(slide, img("ai-translation.png"),
               Inches(7), Inches(1.0), width=Inches(5.8))


# ── 12. HAIQU: RESULTS ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "AI Agent Results", font_size=28, bold=True, color=NAVY, font_name="Libre Baskerville")

tf = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.6),
                 "", font_size=14, italic=True, color=MUTED)
tf.paragraphs[0].text = (
    '"Make real qubits available to users in society... such that those '
    'users can experience quantum computing first hand."')
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = MUTED

# Flow diagram as boxes
flow_items = [
    ("Human\nasks questions", NAVY),
    ("→", BLUE),
    ("AI Agent\nbuilds circuits, runs jobs", BLUE),
    ("→", BLUE),
    ("3 Quantum Chips\nTuna-9 · Garnet · Torino", NAVY),
]
x = Inches(1.5)
for label, col in flow_items:
    if label == "→":
        add_textbox(slide, x, Inches(2.0), Inches(0.5), Inches(0.8),
                    "→", font_size=28, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
        x += Inches(0.6)
    else:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.8), Inches(2.8), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        add_textbox(slide, x + Inches(0.1), Inches(1.9), Inches(2.6), Inches(0.8),
                    label, font_size=12, color=NAVY, alignment=PP_ALIGN.CENTER)
        x += Inches(3.2)

# Stats
stats = [
    ("93%", "claims replicated (25/27)"),
    ("6", "landmark papers"),
    ("0", "circuits by human"),
    ("230K+", "measurement shots"),
]
for i, (val, label) in enumerate(stats):
    cx = Inches(1.2 + i * 3.0)
    add_textbox(slide, cx, Inches(3.6), Inches(2.8), Inches(1),
                val, font_size=40, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx, Inches(4.8), Inches(2.8), Inches(0.5),
                label, font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── 16. SCORECARD ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "Deliverables Scorecard", font_size=28, bold=True, color=NAVY, font_name="Libre Baskerville")

sc_data = [
    ('"PhD thesis (month 48)"', "DONE", "QCT thesis — 6 RQs, framework, metaphors, prototype (PhD)"),
    ('"3 scientific publications"', "DONE", "Thesis chapters (PhD) + AI replication paper (2026)"),
    ('"Fundamental research on mental models"', "DONE", "QCT framework + 39 metaphors + 7 designed (PhD); MSc thesis"),
    ('"Technology probes... co-creation"', "DONE", "Educator interviews + prototype eval n=10 (PhD); Quantum Buddy (MSc)"),
    ('"Updated UX and interface"', "DONE", "AI agent interface + interactive dashboard (haiqu.org)"),
    ('"Data through public repositories"', "DONE", "PhD data online (PhD); open-source repo, 98 results, 230K+ shots"),
    ('"Interfacing to stakeholders" (WP5-7)', "ONGOING", "AI agent approach opens new pathways"),
]

tbl = slide.shapes.add_table(len(sc_data)+1, 3, Inches(0.6), Inches(1.0), Inches(12), Inches(5))
table = tbl.table
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(7.3)

for ci, txt in enumerate(["Grant promise", "Status", "Delivered"]):
    cell = table.cell(0, ci)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY

for ri, (promise, status, delivered) in enumerate(sc_data):
    vals = [promise, status, delivered]
    for ci, val in enumerate(vals):
        cell = table.cell(ri+1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = GREEN if (ci == 1 and val == "DONE") else AMBER if (ci == 1 and val == "ONGOING") else SLATE
            if ci == 1:
                p.font.bold = True
        if ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE


# ── 15. WHAT REMAINS ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
            "What Remains", font_size=28, bold=True, color=NAVY, font_name="Libre Baskerville")

tf = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.6),
                 "", font_size=14, italic=True, color=MUTED)
tf.paragraphs[0].text = (
    '"It is essential to find this common ground between the experts '
    'developing the technology and prospective users."')
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = MUTED

# Left: open questions
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.3),
            "OPEN QUESTIONS", font_size=11, bold=True, color=BLUE)
questions = [
    "•  Does AI abstraction enable learning or obscure understanding?",
    "•  Can QCT metaphors be woven into AI explanations?",
    "•  If AI can do replication science, can experts + AI do greenfield science?",
]
tf = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2),
                 "", font_size=13, color=SLATE)
for i, q in enumerate(questions):
    if i == 0:
        tf.paragraphs[0].text = q
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = SLATE
    else:
        add_para(tf, q, font_size=13, color=SLATE, spacing_before=Pt(6))

# Right: where we need help
add_textbox(slide, Inches(7), Inches(1.8), Inches(5), Inches(0.3),
            "WHERE WE NEED YOUR HELP", font_size=11, bold=True, color=BLUE)
help_items = [
    "•  Evaluation partners: test haiqu.org with your students, researchers, or stakeholders",
    "•  Domain collaboration: bring your experiments to the AI agent",
    "•  WP5-7 integration: connect the AI interface to real stakeholder needs",
]
tf2 = add_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2),
                  "", font_size=13, color=SLATE)
for i, s in enumerate(help_items):
    if i == 0:
        tf2.paragraphs[0].text = s
        tf2.paragraphs[0].font.size = Pt(13)
        tf2.paragraphs[0].font.color.rgb = SLATE
    else:
        add_para(tf2, s, font_size=13, color=SLATE, spacing_before=Pt(6))

# Try it box
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7), Inches(3.8), Inches(5.5), Inches(0.6))
box2.fill.solid()
box2.fill.fore_color.rgb = SURFACE
box2.line.color.rgb = BLUE
box2.line.width = Pt(1)
add_textbox(slide, Inches(7.2), Inches(3.85), Inches(5), Inches(0.5),
            "Try it now at haiqu.org — run a quantum experiment in natural language.",
            font_size=12, bold=True, color=SLATE)

# Core tension box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
box.line.color.rgb = AMBER
box.line.width = Pt(1)
add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
            "Core tension: An AI that fully abstracts QC is useful. But does it educate? Finding the interface that does both is the remaining challenge.",
            font_size=13, bold=True, color=SLATE)


# ── 18. CLOSE ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, NAVY)
add_bg_image(slide, img("title-bg.png"), opacity=0.25)

tf = add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1),
            "From Mental Models to",
            font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
            font_name="Libre Baskerville")
add_para(tf, "AI-Mediated Quantum Computing",
         font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
         font_name="Libre Baskerville", spacing_before=Pt(0))

lines = [
    "The PhD delivered frameworks, metaphors, and prototypes.",
    "The MSc showed ChatGPT as an accessibility layer.",
    "The AI agent proved it works on real hardware.",
    "",
    "Next: connect these streams so AI-mediated quantum computing",
    "is not just powerful, but educational.",
]
tf = add_textbox(slide, Inches(1.5), Inches(2.8), Inches(9), Inches(2.5),
                 "", font_size=16, color=RGBColor(0xCC, 0xCC, 0xDD),
                 alignment=PP_ALIGN.LEFT)
for i, line in enumerate(lines):
    if i == 0:
        tf.paragraphs[0].text = line
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    else:
        add_para(tf, line, font_size=16,
                 color=RGBColor(0xCC, 0xCC, 0xDD),
                 alignment=PP_ALIGN.LEFT)

contact = [
    "J. Derek Lomas — j.d.lomas@tudelft.nl",
    "Caiseal Beardow — c.p.beardow@tudelft.nl",
    "IDE & QuTech, TU Delft · haiqu.org",
]
tf2 = add_textbox(slide, Inches(1.5), Inches(5.5), Inches(9), Inches(1.5),
                  "", font_size=12, color=RGBColor(0x88, 0x99, 0xAA),
                  alignment=PP_ALIGN.LEFT)
for i, c in enumerate(contact):
    if i == 0:
        tf2.paragraphs[0].text = c
        tf2.paragraphs[0].font.size = Pt(12)
        tf2.paragraphs[0].font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
        tf2.paragraphs[0].alignment = PP_ALIGN.LEFT
    else:
        add_para(tf2, c, font_size=12,
                 color=RGBColor(0x88, 0x99, 0xAA),
                 alignment=PP_ALIGN.LEFT)


# ── SAVE ──────────────────────────────────────────────────────
out = os.path.join(DIR, "wp4-4.pptx")
prs.save(out)
print(f"Saved: {out}")
